import sys
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import os
os.chdir("D:\\revere\\office\\uploads");
print("Output from Python PPt")
mylist = sys.argv[1].split(',')
destination = sys.argv[2]
mainPrs = Presentation(destination)
blank_slide_layout = mainPrs.slide_layouts[5]
left = top = Inches(0)
for a in mylist:
    print(a)
    prs = Presentation(a);
    s=0
    for slide in prs.slides:
        a=0
        b=0
        text_runs = ""
        if(len(slide.shapes)>0):
            mainSlide = mainPrs.slides.add_slide(blank_slide_layout)
            for shape in slide.shapes:
                print(shape.shape_type)
                if(shape.is_placeholder):
                    h= shape.placeholder_format.type
                    print(h)
                    try:
                        if(h == 18):
                            path = str(s)+'mypic.jpg'
                            with open(path, 'wb') as f:
                                f.write(shape.image.blob)
                                s=s+1
                            try:
                                mainSlide.shapes.add_picture(path, left, top, height = mainPrs.slide_height)
                            except:
                                print("Unable to add to slide")
                            os.remove(path)
                    except:
                      print("An exception occurred")
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs = text_runs + run.text
                if(shape.shape_type == 13):
                    a=a+1
                    path = str(s)+'mypic.jpg'
                    with open(path, 'wb') as f:
                        f.write(shape.image.blob)
                        s=s+1
                    try:
                        mainSlide.shapes.add_picture(path, left, top, height = mainPrs.slide_height)
                    except:
                        print("Unable to add to slide")
                    os.remove(path)
                if(shape.shape_type == 1):
                    b=b+1
                if(shape.shape_type == 6):
                    for shape in shape.shapes:
                        print(shape.shape_type)
                        if(shape.is_placeholder):
                            h= shape.placeholder_format.type
                            print(h)
                            try:
                                if(h == 18):
                                    path = str(s)+'mypic.jpg'
                                    with open(path, 'wb') as f:
                                        f.write(shape.image.blob)
                                        s=s+1
                                    try:
                                        mainSlide.shapes.add_picture(path, left, top, height = mainPrs.slide_height)
                                    except:
                                        print("Unable to add to slide")
                                    os.remove(path)
                            except:
                              print("An exception occurred")
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text_runs = text_runs + run.text
                        if(shape.shape_type == 13):
                            a=a+1
                            path = str(s)+'mypic.jpg'
                            with open(path, 'wb') as f:
                                f.write(shape.image.blob)
                                s=s+1
                            try:
                                mainSlide.shapes.add_picture(path, left, top, height = mainPrs.slide_height)
                            except:
                                print("Unable to add to slide")
                            os.remove(path)
            print("a= "+str(a)+",b= "+str(b))
        halfI = Inches(0.5)
        textBoxTop = mainPrs.slide_height - halfI
        textBox = mainSlide.shapes.add_textbox(left, top = textBoxTop,width = mainPrs.slide_width, height = halfI)
        fill = textBox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(235, 192, 52)
        textBox.text_frame.word_wrap = True
        p = textBox.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text_runs
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(24)
        font.bold = True
print("print outside loop before saving")
mainPrs.save(destination)
print("file saved to " + destination)
